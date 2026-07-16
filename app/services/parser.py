import pdfplumber
import docx
import pandas as pd
from pptx import Presentation
from pathlib import Path


def _read_text_file(file_path: str) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
        try:
            with open(file_path, encoding=encoding) as file:
                text = file.read()
                if text.strip():
                    return text
        except UnicodeDecodeError:
            continue
    raise ValueError("No text extracted")


def extract_text(file_path: str) -> str:
    
    suffix = Path(file_path).suffix

    if suffix == '.pdf':
        with pdfplumber.open(file_path) as pdf:
            pages_text = [text for page in pdf.pages if (text := page.extract_text())]
            if pages_text:
                return '\n'.join(pages_text)
            else:
                raise ValueError('No text extracted')
    
    elif suffix in ('.md', '.txt'):
        return _read_text_file(file_path)

    elif suffix == '.docx':
        doc = docx.Document(file_path)
        paragraphs = [text for p in doc.paragraphs if (text := p.text.strip())]
        if paragraphs:
            return '\n'.join(paragraphs)
        else:
            raise ValueError('No text extracted')
    
    elif suffix == '.xlsx':
        df_dict = pd.read_excel(file_path, sheet_name=None)
        text = [f'{sheet_name}: {df.to_string(index=False).strip()}' 
                for sheet_name, df in df_dict.items() 
                if not df.empty
        ]
        if text:
            return '\n\n'.join(text)
        else:
            raise ValueError('No text extracted')

    elif suffix == '.pptx':
        presentation = Presentation(file_path)
        slides = [
            text
            for slide in presentation.slides
            for shape in slide.shapes 
            if hasattr(shape, 'text') and (text := shape.text)
            ]
        if slides:
            return '\n'.join(slides)
        else:
            raise ValueError('No text extracted')
    
    else:
        raise ValueError('Unsupported file type')
    